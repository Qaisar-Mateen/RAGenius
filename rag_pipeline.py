import os
import logging
import time
import pickle
from typing import List, Dict, Tuple, Optional, Any
from dotenv import load_dotenv

# Vector DB imports
from llama_index.vector_stores.qdrant import QdrantVectorStore
from llama_index.core import VectorStoreIndex, Document, StorageContext, Settings
from llama_index.core.node_parser import SentenceSplitter
import qdrant_client

# Embedding model
from llama_index.embeddings.fastembed import FastEmbedEmbedding

# LLM
from llama_index.llms.groq import Groq

# Document processing
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from llama_parse import LlamaParse  # Added LlamaParse

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class RAGPipeline:
    """RAG Pipeline for document processing and retrieval."""
    
    def __init__(        self, 
        qdrant_url: Optional[str] = None, 
        qdrant_api_key: Optional[str] = None,
        groq_api_key: Optional[str] = None,
        llamaparse_api_key: Optional[str] = None,  # Added LlamaParse API key
        collection_name: str = "study_materials",
        embedding_model: str = "BAAI/bge-base-en-v1.5",
        llm_model: str = "llama-3.3-70b-versatile",
        storage_dir: str = "./storage",
        chunk_size: int = 1500,
        chunk_overlap: int = 200
    ):
        self.qdrant_url = qdrant_url or os.getenv("QDRANT_URL")
        self.qdrant_api_key = qdrant_api_key or os.getenv("QDRANT_API_KEY")
        self.groq_api_key = groq_api_key or os.getenv("GROQ_API_KEY")
        self.llamaparse_api_key = llamaparse_api_key or os.getenv("PARSE_API_KEY")
        self.collection_name = collection_name
        self.embedding_model_name = embedding_model
        self.llm_model_name = llm_model
        self.storage_dir = storage_dir
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        
        # Create storage directory if it doesn't exist
        os.makedirs(self.storage_dir, exist_ok=True)
        
        # Ensure temp directory exists for document processing
        os.makedirs("./data", exist_ok=True)
        
        # Initialize paragraph-aware node parser for better text splitting
        from llama_index.core.node_parser import SentenceSplitter
        from llama_index.core.text_splitter import TokenTextSplitter

        # Create a node parser that properly chunks documents by paragraphs
        # Using SentenceSplitter with paragraph separator configuration
        self.node_parser = SentenceSplitter(
            separator="\n\n",  # Empty line indicates paragraph break
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap
        )
        
        # Initialize components
        self._initialize_embedding_model()
        self._initialize_llm()
        
        # Set up vector store and other components
        self.qdrant_client = None
        self.vector_store = None
        self.index = None
        self.query_engine = None
        
        # Initialize vector store
        self._initialize_vector_store()

    def _initialize_embedding_model(self):
        """Initialize the embedding model."""
        try:
            self.embed_model = FastEmbedEmbedding(model_name=self.embedding_model_name)
            Settings.embed_model = self.embed_model
            logger.info(f"Embedding model initialized: {self.embedding_model_name}")
        except Exception as e:
            logger.error(f"Error initializing embedding model: {e}")
            raise

    def _initialize_llm(self):
        """Initialize the language model."""
        try:
            if not self.groq_api_key:
                logger.warning("Groq API key not provided. LLM not initialized.")
                self.llm = None
                return
                
            self.llm = Groq(model=self.llm_model_name, api_key=self.groq_api_key)
            Settings.llm = self.llm
            logger.info(f"Language model initialized: {self.llm_model_name}")
        except Exception as e:
            logger.error(f"Error initializing language model: {e}")
            raise

    def _initialize_vector_store(self):
        """Initialize the vector store."""
        if not self.qdrant_url or not self.qdrant_api_key:
            logger.warning("Qdrant URL or API key not provided. Vector store not initialized.")
            return
            
        try:
            self.qdrant_client = qdrant_client.QdrantClient(
                api_key=self.qdrant_api_key, 
                url=self.qdrant_url
            )
            
            self.vector_store = QdrantVectorStore(
                client=self.qdrant_client, 
                collection_name=self.collection_name
            )
            
            logger.info(f"Vector store initialized: Qdrant ({self.collection_name})")
        except Exception as e:
            logger.error(f"Error initializing vector store: {e}")
            raise
            
    def clear_collection(self) -> bool:
        """Clear the Qdrant collection to remove old data before indexing new files.
        
        Returns:
            bool: True if successful, False otherwise
        """
        if not self.qdrant_client:
            logger.warning("Qdrant client not initialized. Cannot clear collection.")
            return False
            
        try:
            # Check if collection exists
            collections = self.qdrant_client.get_collections().collections
            collection_names = [col.name for col in collections]
            
            if self.collection_name in collection_names:
                # Delete collection
                logger.info(f"Deleting existing collection: {self.collection_name}")
                self.qdrant_client.delete_collection(collection_name=self.collection_name)
                
                # Recreate collection with same settings
                logger.info(f"Recreating collection: {self.collection_name}")
                dimension = 768  # Default dimension for BAAI/bge-base-en-v1.5
                self.qdrant_client.create_collection(
                    collection_name=self.collection_name,
                    vectors_config=qdrant_client.http.models.VectorParams(
                        size=dimension,
                        distance=qdrant_client.http.models.Distance.COSINE
                    )
                )
                logger.info(f"Collection {self.collection_name} cleared and recreated successfully")
                
                # Reinitialize vector store with new collection
                self.vector_store = QdrantVectorStore(
                    client=self.qdrant_client, 
                    collection_name=self.collection_name
                )
                return True
            else:
                logger.info(f"Collection {self.collection_name} does not exist yet, nothing to clear")
                return True
        except Exception as e:
            logger.error(f"Error clearing collection: {e}")
            return False

    def extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF file."""
        try:
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting text from PDF {file_path}: {e}")
            raise

    def extract_text_from_docx(self, file_path: str) -> str:
        """Extract text from Word document."""
        try:
            doc = DocxDocument(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                if paragraph.text:
                    text += paragraph.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting text from DOCX {file_path}: {e}")
            raise

    def extract_text_from_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint presentation."""
        try:
            ppt = Presentation(file_path)
            text = ""
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting text from PPTX {file_path}: {e}")
            raise

    def process_uploaded_file(self, uploaded_file, temp_dir: str) -> Tuple[Optional[str], Optional[str]]:
        """Process an uploaded file and extract text."""
        file_path = os.path.join(temp_dir, uploaded_file.name)
        
        try:
            # Save uploaded file to disk temporarily
            with open(file_path, "wb") as f:
                f.write(uploaded_file.getbuffer())
            
            # Extract text based on file extension
            file_extension = os.path.splitext(uploaded_file.name)[1].lower()
            
            if file_extension == '.pdf':
                text = self.extract_text_from_pdf(file_path)
            elif file_extension == '.docx':
                text = self.extract_text_from_docx(file_path)
            elif file_extension in ['.pptx', '.ppt']:
                text = self.extract_text_from_pptx(file_path)
            else:
                return None, f"Unsupported file format: {file_extension}"
            
            # Remove the temporary file
            os.remove(file_path)
            
            # Check if we actually got content
            if not text or len(text.strip()) < 10:
                return None, f"Could not extract meaningful text from {uploaded_file.name}"
                
            return text, None
        except Exception as e:
            # Remove the temporary file if it exists
            if os.path.exists(file_path):
                os.remove(file_path)
            return None, f"Error processing {uploaded_file.name}: {str(e)}"

    def process_documents(self, documents_dict: Dict[str, str]) -> List[Document]:
        """Convert extracted text to LlamaIndex documents with chunking."""
        all_nodes = []
        
        for filename, text in documents_dict.items():
            # Create a document
            doc = Document(
                text=text,
                metadata={"filename": filename, "source": filename}
            )
            
            # Parse the document into chunks
            nodes = self.node_parser.get_nodes_from_documents([doc])
            
            # Ensure each node has the source filename
            for node in nodes:
                if "filename" not in node.metadata:
                    node.metadata["filename"] = filename
            
            all_nodes.extend(nodes)
            
        logger.info(f"Created {len(all_nodes)} chunks from {len(documents_dict)} documents")
        
        # Log chunks to a text file
        log_path = self._log_chunks_to_file(all_nodes, f"chunks_log_{int(time.time())}.txt")
        if log_path:
            logger.info(f"Chunks saved to: {log_path}")
        
        return all_nodes

    def build_index(self, documents: List[Document], clear_database: bool = True) -> VectorStoreIndex:
        """Build vector index from documents.
        
        Args:
            documents: List of documents to index
            clear_database: If True, clear the Qdrant database before indexing
            
        Returns:
            VectorStoreIndex: The built index
        """
        try:
            if not self.vector_store:
                logger.error("Vector store not initialized. Cannot build index.")
                raise ValueError("Vector store not initialized")
            
            # Clear the Qdrant collection before indexing if requested
            if clear_database:
                logger.info("Clearing Qdrant collection before building new index")
                success = self.clear_collection()
                if not success:
                    logger.error("Failed to clear Qdrant collection")
                    raise ValueError("Failed to clear Qdrant collection")
                
            storage_context = StorageContext.from_defaults(vector_store=self.vector_store)
            
            index = VectorStoreIndex(
                nodes=documents, 
                storage_context=storage_context,
                show_progress=True
            )
            
            self.index = index
            return index
        except Exception as e:
            logger.error(f"Error building index: {e}")
            raise

    def save_index(self, index_name: str = "study_materials_index"):
        """Save the index to disk."""
        if not self.index:
            logger.warning("No index to save")
            return False
            
        try:
            # Persist the index in the storage directory
            index_path = os.path.join(self.storage_dir, index_name)
            
            # Clean the existing directory if it exists to avoid conflicts
            if os.path.exists(index_path):
                # This helps ensure we don't have orphaned files from previous indices
                logger.info(f"Cleaning existing index directory: {index_path}")
                for filename in os.listdir(index_path):
                    file_path = os.path.join(index_path, filename)
                    try:
                        if os.path.isfile(file_path):
                            os.unlink(file_path)
                    except Exception as e:
                        logger.warning(f"Error deleting {file_path}: {e}")
            
            # Make sure directory exists
            os.makedirs(index_path, exist_ok=True)
            
            # Persist the index
            self.index.storage_context.persist(persist_dir=index_path)
            
            logger.info(f"Index saved to {index_path}")
            return True
        except Exception as e:
            logger.error(f"Error saving index: {e}")
            return False

    def create_query_engine(self, similarity_top_k=7):
        """Create a query engine from the index with configurable parameters."""
        if not self.index:
            logger.warning("No index available for query engine")
            return None
            
        self.query_engine = self.index.as_query_engine(
            similarity_top_k=similarity_top_k,
            llm=self.llm if self.llm else None
        )
        return self.query_engine

    def query(self, query_text: str, similarity_top_k=7) -> Dict[str, Any]:
        """Query the RAG system with configurable parameters."""
        start_time = time.time()
        
        if not self.index:
            return {
                "answer": "Sorry, no documents have been indexed yet. Please upload study materials first.",
                "sources": [],
                "elapsed_time": 0
            }
        
        try:
            # Create or update query engine with the specified parameters
            self.create_query_engine(similarity_top_k=similarity_top_k)
                
            # Execute query
            response = self.query_engine.query(query_text)
            
            # Extract sources and prepare response
            sources = []
            if hasattr(response, 'source_nodes'):
                sources = [
                    {
                        "filename": node.metadata.get("filename", "Unknown"),
                        "text": node.text[:200] + "..." if len(node.text) > 200 else node.text,
                        "score": node.score if hasattr(node, 'score') else None
                    }
                    for node in response.source_nodes
                ]
            
            elapsed_time = time.time() - start_time
            
            return {
                "answer": str(response),
                "sources": sources,
                "elapsed_time": elapsed_time
            }
        except Exception as e:
            logger.error(f"Error querying index: {e}")
            return {
                "answer": f"An error occurred while processing your query: {str(e)}",
                "sources": [],
                "elapsed_time": time.time() - start_time
            }

    def process_and_index_files(self, uploaded_files_info: List[Dict]) -> bool:
        """Process and index files from uploaded_files_info."""
        try:
            # Extract file paths from file info for LlamaParse processing
            file_paths = []
            temp_files = []
            file_names_map = {}  # Maps temp file paths to original filenames for metadata
            
            # First, save uploaded files to disk temporarily
            for file_info in uploaded_files_info:
                if "name" in file_info:
                    # Create temp file path
                    temp_path = os.path.join("./data", file_info["name"])
                    temp_files.append(temp_path)
                    
                    # Write content to temp file if available
                    if "content" in file_info and file_info["content"]:
                        with open(temp_path, "wb" if isinstance(file_info["content"], bytes) else "w", encoding=None if isinstance(file_info["content"], bytes) else "utf-8") as f:
                            f.write(file_info["content"])
                        file_paths.append(temp_path)
                        file_names_map[temp_path] = file_info["name"]
            
            if not file_paths:
                logger.warning("No valid files found for processing")
                return False
            
            # Use LlamaParse to process the documents
            documents = None
            try:
                # Try LlamaParse first as the preferred method
                logger.info(f"Attempting to parse {len(file_paths)} files with LlamaParse")
                documents = self.load_or_parse_data(file_paths)
                
                # Ensure metadata preservation - sometimes LlamaParse loses original filenames
                for doc in documents:
                    if "source" in doc.metadata and doc.metadata["source"] in file_names_map:
                        doc.metadata["filename"] = file_names_map[doc.metadata["source"]]
                    elif "source_path" in doc.metadata and doc.metadata["source_path"] in file_names_map:
                        doc.metadata["filename"] = file_names_map[doc.metadata["source_path"]]
                
                logger.info(f"Successfully parsed {len(documents)} documents using LlamaParse")
                
                # Log the chunks from LlamaParse processing
                log_path = self._log_chunks_to_file(documents, f"llama_chunks_log_{int(time.time())}.txt")
                if log_path:
                    logger.info(f"LlamaParse chunks saved to: {log_path}")
            except Exception as e:
                logger.warning(f"LlamaParse parsing failed: {e}. Falling back to traditional parsing.")
                documents = None
                
            # Fallback to traditional document parsing if LlamaParse failed
            if not documents:
                logger.info("Using traditional document parsing methods")
                documents_dict = {}
                for file_path in file_paths:
                    file_name = file_names_map.get(file_path, os.path.basename(file_path))
                    file_extension = os.path.splitext(file_name)[1].lower()
                    
                    try:
                        if file_extension == '.pdf':
                            text = self.extract_text_from_pdf(file_path)
                        elif file_extension == '.docx':
                            text = self.extract_text_from_docx(file_path)
                        elif file_extension in ['.pptx', '.ppt']:
                            text = self.extract_text_from_pptx(file_path)
                        else:
                            logger.warning(f"Unsupported file format: {file_extension}")
                            continue
                            
                        documents_dict[file_name] = text
                    except Exception as extract_err:
                        logger.warning(f"Error extracting text from {file_name}: {extract_err}")
                
                if not documents_dict:
                    logger.warning("No content found in uploaded files")
                    return False
                
                # Convert to chunked LlamaIndex documents
                documents = self.process_documents(documents_dict)
            
            # Clean up temporary files
            for temp_file in temp_files:
                if os.path.exists(temp_file):
                    os.remove(temp_file)
            
            if not documents:
                logger.warning("No documents were parsed successfully")
                return False
                
            # Build the index
            self.build_index(documents)
            
            # Save the index
            self.save_index()
            
            # Create query engine
            self.create_query_engine()
            
            return True
        except Exception as e:
            logger.error(f"Error processing and indexing files: {e}", exc_info=True)
            # Clean up any temporary files
            for temp_file in temp_files if 'temp_files' in locals() else []:
                if os.path.exists(temp_file):
                    try:
                        os.remove(temp_file)
                    except Exception:
                        pass
            return False
            
    def test_pipeline(self) -> Dict[str, Any]:
        """Run a simple test to verify the pipeline is working."""
        if not self.index:
            return {
                "status": "error",
                "message": "No index available. Please upload and index documents first."
            }
            
        try:
            # Try a simple query
            test_query = "What is this document about?"
            result = self.query(test_query, similarity_top_k=1)
            
            if result and "answer" in result and result["answer"]:
                return {
                    "status": "success",
                    "message": "RAG pipeline is working correctly.",
                    "test_query": test_query,
                    "result": result
                }
            else:
                return {
                    "status": "warning",
                    "message": "Query returned empty results. Check document content."
                }
        except Exception as e:
            return {
                "status": "error",
                "message": f"Error testing pipeline: {str(e)}"
            }

    def load_or_parse_data(self, file_paths: List[str], cache_dir: str = "./data") -> List[Document]:
        """Load parsed data if available, or parse using LlamaParse if not.
        
        Args:
            file_paths: List of file paths to parse
            cache_dir: Directory to store cached parsed data
            
        Returns:
            List of LlamaIndex Document objects
        """
        if not file_paths:
            logger.warning("No file paths provided for parsing")
            return []
            
        # Create a session-specific cache ID based on current timestamp
        # This ensures we only use documents from the current chat session
        session_id = int(time.time())
        
        # Create unique cache file name based on the file paths and session ID
        # Including session_id ensures we never reuse cached data across sessions
        files_hash = "_".join([os.path.basename(f) for f in file_paths])
        cache_filename = f"llama_parsed_{session_id}_{files_hash}.pkl"
        cache_filename = cache_filename.replace(' ', '_')[:100]  # Limit filename length
        cache_file = os.path.join(cache_dir, cache_filename)
        
        # We'll always parse fresh for each session to ensure we only use current documents
        
        # Check if LlamaParse API key is available
        if not self.llamaparse_api_key:
            logger.error("LlamaParse API key not provided. Cannot parse documents.")
            raise ValueError("LlamaParse API key is required for document parsing")
        
        try:
            # Create a mapping of file paths to original filenames for metadata preservation
            file_names_map = {path: os.path.basename(path) for path in file_paths}
            file_basenames = [os.path.basename(path) for path in file_paths]
            
            # Parse documents using LlamaParse
            logger.info(f"Parsing {len(file_paths)} files with LlamaParse")
            llama_parse_documents = LlamaParse(
                api_key=self.llamaparse_api_key, 
                result_type="markdown"
            ).load_data(file_paths)
            
            # Ensure metadata preservation for current session files
            for doc in llama_parse_documents:
                source_path = doc.metadata.get("source_path", "")
                source_filename = os.path.basename(source_path) if source_path else ""
                
                # Check if the document matches any of our input files by basename
                if source_path in file_names_map or source_filename in file_basenames:
                    # Document is from current session
                    doc.metadata["filename"] = file_names_map.get(source_path, source_filename)
                    doc.metadata["original_filename"] = file_names_map.get(source_path, source_filename)
                    doc.metadata["session_id"] = session_id
                else:
                    # More lenient matching - check if source path ends with any of our filenames
                    matched = False
                    for base_name in file_basenames:
                        if source_path.endswith(base_name):
                            doc.metadata["filename"] = base_name
                            doc.metadata["original_filename"] = base_name
                            doc.metadata["session_id"] = session_id
                            matched = True
                            break
                            
                    if not matched:
                        # If file wasn't in our list, add a flag but don't filter yet
                        logger.warning(f"Document found with source path not in current session: {source_path}")
                        # We'll keep all documents for now 
                        # doc.metadata["not_current_session"] = True
            
            # ** DISABLED FILTERING ** - Use all documents returned by LlamaParse for now
            # Filter out any documents not from the current session
            # current_session_docs = [doc for doc in llama_parse_documents 
            #                        if "not_current_session" not in doc.metadata]
            current_session_docs = llama_parse_documents  # Use all documents
            
            if len(current_session_docs) < len(llama_parse_documents):
                logger.info(f"Filtered out {len(llama_parse_documents) - len(current_session_docs)} documents not from current session")
            
            # Apply paragraph-wise chunking to the LlamaParse documents
            logger.info(f"Applying paragraph-wise chunking to {len(current_session_docs)} LlamaParse documents")
            chunked_nodes = []
            for doc in current_session_docs:
                # Apply the node_parser to get paragraph-wise chunks
                nodes = self.node_parser.get_nodes_from_documents([doc])
                
                # Preserve metadata from the original document
                for node in nodes:
                    node.metadata.update(doc.metadata)
                
                chunked_nodes.extend(nodes)
            
            logger.info(f"Created {len(chunked_nodes)} paragraph-wise chunks from {len(current_session_docs)} LlamaParse documents")
            
            # Save parsed and chunked data to cache
            os.makedirs(cache_dir, exist_ok=True)
            with open(cache_file, "wb") as f:
                pickle.dump(chunked_nodes, f)
            
            logger.info(f"Parsed documents saved to {cache_file}")
            
            return chunked_nodes
        except Exception as e:
            logger.error(f"Error parsing documents with LlamaParse: {e}")
            raise

    def _log_chunks_to_file(self, nodes, filename="chunks_log.txt"):
        """Log chunked data to a text file for inspection.
        
        Args:
            nodes: List of nodes (chunks) to log
            filename: Name of the log file
            
        Returns:
            str: Path to the log file or None if an error occurred
        """
        log_dir = os.path.join(self.storage_dir, "logs")
        os.makedirs(log_dir, exist_ok=True)
        
        # Create timestamp for unique filename
        timestamp = time.strftime("%Y%m%d_%H%M%S")
        log_filename = f"{timestamp}_{filename}"
        log_path = os.path.join(log_dir, log_filename)
        
        try:
            with open(log_path, "w", encoding="utf-8") as f:
                f.write(f"Chunked Data Log - Created: {time.strftime('%Y-%m-%d %H:%M:%S')}\n")
                f.write(f"Total Chunks: {len(nodes)}\n")
                f.write(f"Chunk Size: {self.chunk_size}, Chunk Overlap: {self.chunk_overlap}\n\n")
                
                for i, node in enumerate(nodes):
                    f.write(f"{'='*80}\n")
                    f.write(f"CHUNK #{i+1}\n")
                    f.write(f"Source: {node.metadata.get('filename', 'Unknown')}\n")
                    f.write(f"Node ID: {node.node_id}\n")
                    f.write(f"{'='*80}\n\n")
                    f.write(node.text)
                    f.write("\n\n")
            
            logger.info(f"Chunked data logged to {log_path}")
            return log_path
        except Exception as e:
            logger.error(f"Error logging chunks to file: {e}")
            return None