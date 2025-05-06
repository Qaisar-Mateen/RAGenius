from typing import Generator, List, Dict, Any
import re
import time
import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter

def icon(emoji: str, st):
    """Shows an emoji as a Notion-style page icon."""
    st.write(
        f'<span style="font-size: 78px; line-height: 1">{emoji}</span>',
        unsafe_allow_html=True,
    )

def extract_thinking(content: str) -> tuple:
    """Extract thinking part and actual content from a message."""
    thinking_pattern = r'<think>(.*?)</think>'
    thinking_match = re.search(thinking_pattern, content, re.DOTALL)
    
    if thinking_match:
        thinking = thinking_match.group(1).strip()
        # Remove thinking part from the content
        actual_content = re.sub(thinking_pattern, '', content, flags=re.DOTALL).strip()
        return thinking, actual_content
    return None, content

def remove_thinking(content: str) -> str:
    """Remove thinking part from content."""
    _, clean_content = extract_thinking(content)
    return clean_content

def stream_Chat(chat_completion) -> Generator[tuple, None, None]:
    """Yield chat response content from the Groq API response with thinking parts separated."""
    full_text = ""
    current_thinking = None
    previous_thinking = None
    current_content = ""
    previous_content = ""
    in_thinking_section = False
    
    for chunk in chat_completion:
        if not chunk.choices[0].delta.content:
            continue
            
        # Get the current chunk and add it to full text
        chunk_text = chunk.choices[0].delta.content
        full_text += chunk_text
        
        # More efficient handling of thinking sections
        think_start = full_text.find("<think>") 
        think_end = full_text.find("</think>")
        
        # Complete thinking section found
        if think_start >= 0 and think_end > think_start:
            in_thinking_section = False
            current_thinking = full_text[think_start+7:think_end].strip()
            
            # Get content before and after thinking section
            before_thinking = full_text[:think_start].strip()
            after_thinking = full_text[think_end+8:].strip()
            current_content = before_thinking + " " + after_thinking if before_thinking and after_thinking else before_thinking + after_thinking
            
            # Remove processed thinking section from full_text
            full_text = full_text[:think_start] + full_text[think_end+8:]
            
        # Partial thinking section (started but not completed)
        elif think_start >= 0:
            in_thinking_section = True
            current_thinking = full_text[think_start+7:].strip()
            current_content = full_text[:think_start].strip()
        # No thinking section
        elif not in_thinking_section:
            current_content = full_text.strip()
            current_thinking = None
        
        # Add a small delay for slower generation
        time.sleep(0.01)
        
        # Only yield if content has changed to avoid redundant updates
        if (current_thinking != previous_thinking) or (current_content != previous_content):
            previous_thinking = current_thinking
            previous_content = current_content
            
            yield {
                "thinking": current_thinking,
                "content": current_content,
                "chunk": chunk_text
            }

def extract_text_from_pdf(file_path):
    """Extract text from PDF file with improved paragraph structure."""
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            # Clean up PDF text to better preserve paragraphs
            # Replace multiple newlines with paragraph breaks
            page_text = re.sub(r'\n\s*\n', '\n\n', page_text)
            # Replace single newlines that don't end with sentence-ending punctuation
            # with spaces (preserving true paragraph breaks)
            page_text = re.sub(r'([^\.\?\!])\n([A-Z0-9])', r'\1 \2', page_text)
            text += page_text + "\n\n"  # Add paragraph break between pages
    return text

def extract_text_from_docx(file_path):
    """Extract text from Word document with improved paragraph structure."""
    doc = Document(file_path)
    paragraphs = []
    
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            # Add non-empty paragraphs
            paragraphs.append(paragraph.text.strip())
    
    # Join paragraphs with double newlines to preserve paragraph structure
    return "\n\n".join(paragraphs)

def extract_text_from_pptx(file_path):
    """Extract text from PowerPoint presentation with slide separation."""
    ppt = Presentation(file_path)
    slides_text = []
    
    for i, slide in enumerate(ppt.slides):
        slide_parts = []
        # Add slide number as header
        slide_parts.append(f"Slide {i+1}:")
        
        # Process shapes in the slide
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
        
        # Join all text from this slide
        if len(slide_parts) > 1:  # Only include if there's content beyond the header
            slides_text.append("\n".join(slide_parts))
    
    # Join slides with double newlines to create clear separation
    return "\n\n".join(slides_text)

def process_uploaded_file(uploaded_file, temp_dir):
    """Process an uploaded file and extract text."""
    # Create a temporary file path
    file_path = os.path.join(temp_dir, uploaded_file.name)
    
    # Save uploaded file to disk temporarily
    with open(file_path, "wb") as f:
        f.write(uploaded_file.getbuffer())
    
    # Extract text based on file extension
    file_extension = os.path.splitext(uploaded_file.name)[1].lower()
    
    try:
        if file_extension == '.pdf':
            text = extract_text_from_pdf(file_path)
        elif file_extension == '.docx':
            text = extract_text_from_docx(file_path)
        elif file_extension in ['.pptx', '.ppt']:
            text = extract_text_from_pptx(file_path)
        else:
            return None, f"Unsupported file format: {file_extension}"
        
        # Remove the temporary file
        os.remove(file_path)
        return text, None
    except Exception as e:
        # Remove the temporary file if it exists
        if os.path.exists(file_path):
            os.remove(file_path)
        return None, f"Error processing {uploaded_file.name}: {str(e)}"

def chunk_text(text, chunk_size=1500, chunk_overlap=150):
    """Split text into chunks for processing with improved paragraph preservation."""
    if not text:
        return []
        
    # Use a paragraph-aware text splitter that better preserves document structure
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
        separators=["\n\n", "\n", ". ", " ", ""]  # Prioritize paragraph breaks first
    )
    
    # Add paragraph indicators to help the model understand structure
    chunks = text_splitter.split_text(text)
    
    # Format chunks to better indicate paragraphs
    formatted_chunks = []
    for chunk in chunks:
        # Replace paragraph breaks with clear indicators
        formatted_chunk = chunk.replace("\n\n", "\n\n[PARAGRAPH]\n\n")
        formatted_chunks.append(formatted_chunk)
    
    return formatted_chunks